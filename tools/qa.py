# -*- coding: utf-8 -*-
"""덱 기하 검사 — python3 tools/qa.py 발표자료_배터리코칭.pptx"""
import sys, math, unicodedata
from pptx import Presentation
prs=Presentation(sys.argv[1])
SW,SH=prs.slide_width/914400, prs.slide_height/914400
wd=lambda c: 1.0 if unicodedata.east_asian_width(c) in ('W','F') else 0.52
rect=lambda sh:(sh.left/914400, sh.top/914400,(sh.width or 0)/914400,(sh.height or 0)/914400)
issues=[]
for si,sl in enumerate(prs.slides,1):
    tb=[]
    for sh in sl.shapes:
        if sh.left is None: continue
        x,y,w,h=rect(sh)
        if x<-0.02 or y<-0.02 or x+w>SW+0.02 or y+h>SH+0.02:
            lbl=sh.text_frame.text.strip()[:20] if sh.has_text_frame else str(sh.shape_type)
            issues.append(f'S{si} 캔버스 이탈 [{lbl}]')
        if not sh.has_text_frame: continue
        tf=sh.text_frame
        if not tf.text.strip(): continue
        if x<0.47 or x+w>SW-0.47:
            issues.append(f'S{si} 여백 부족 "{tf.text.strip()[:22]}"')
        th=0
        for p in tf.paragraphs:
            t=''.join(r.text for r in p.runs)
            if not t: th+=6/72; continue
            fs=next((r.font.size.pt for r in p.runs if r.font.size),12)
            for seg in t.split('\n'):
                u=sum(wd(c) for c in seg) or 0.5
                th+=max(1,math.ceil(u*fs/72/max(w-0.12,0.25)))*fs*1.34/72
        if th>h+0.08:
            issues.append(f'S{si} 넘침 "{tf.text.strip()[:26]}" {th:.2f}">{h:.2f}"')
        tb.append((x,y,w,h,tf.text.strip()))
    for i in range(len(tb)):
        for j in range(i+1,len(tb)):
            ax,ay,aw,ah,at=tb[i]; bx,by,bw,bh,bt=tb[j]
            ix=max(0,min(ax+aw,bx+bw)-max(ax,bx)); iy=max(0,min(ay+ah,by+bh)-max(ay,by))
            ar=ix*iy; small=min(aw*ah,bw*bh)
            if ar>0 and small>0 and ar/small>0.6:
                issues.append(f'S{si} 텍스트 겹침 "{at[:16]}" ↔ "{bt[:16]}"')
print('\n'.join(issues) if issues else '문제 없음')
print(f'--- 슬라이드 {len(prs.slides)}장 · 문제 {len(issues)}건')
